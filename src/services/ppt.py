from typing import BinaryIO, Iterable, List

from pptx import Presentation
from pptx.exc import PackageNotFoundError


ALLOWED_EXTENSIONS = {".pptx"}
ALLOWED_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/octet-stream",
}


def _iter_text(paragraphs: Iterable) -> List[str]:
    lines = []
    for paragraph in paragraphs:
        text = getattr(paragraph, "text", "").strip()
        if text:
            lines.append(text)
    return lines


def extract_text(file_obj: BinaryIO) -> str:
    """Extract visible text from a PPTX file-like object."""
    try:
        file_obj.seek(0)
        presentation = Presentation(file_obj)
    except PackageNotFoundError as exc:  # type: ignore[attr-defined]
        raise ValueError("Uploaded file is not a valid PPTX document.") from exc

    lines: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append(text)
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        segments = _iter_text(cell.text_frame.paragraphs if cell.text_frame else [])
                        lines.extend(segments)
    return "\n".join(lines)
